#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将 SPXOpenAPILandscape 视频的 7 个场景静态帧组装为 .pptx。

运行前提：
  1. 已通过 `npm run render:spxopenapi-landscape:slides` 渲染出
     out/slides/spxopenapi-landscape/slide-{1..7}.png
  2. 已安装 python-pptx：`pip3 install --user python-pptx pillow`

产物：
  out/SPXOpenAPILandscape-slides.pptx

每页结构：
  - 页面比例：16:9（13.333 × 7.5 inch，1920×1080 的 PPT 规格）
  - 场景页：背景铺 PNG 原图 + 底部条带展示场景标题
  - 每页备注区（Notes）：本场景的口播旁白
  - 首页：产品标题页（纯文本）
  - 末页：安装 / 能力总览页（纯文本）
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Tuple

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# 路径配置
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = ROOT / "out" / "slides" / "spxopenapi-landscape"
OUTPUT_PPTX = ROOT / "out" / "SPXOpenAPILandscape-slides.pptx"

# ---------------------------------------------------------------------------
# 色板（与 schema.ts 保持一致）
# ---------------------------------------------------------------------------
BG = RGBColor(0x07, 0x0A, 0x10)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)   # cyan
HIGHLIGHT = RGBColor(0x8B, 0x5C, 0xF6)  # purple
GOLD = RGBColor(0xFF, 0xD7, 0x00)

# ---------------------------------------------------------------------------
# 场景元数据（标题 + Key Points + 旁白）——与 schema.ts voiceoverScripts 对齐
# ---------------------------------------------------------------------------
Scene = Tuple[str, str, List[str], str]  # (编号, 标题, 要点列表, 旁白)

SCENES: List[Scene] = [
    (
        "Scene 1 · HOOK",
        "物流 API 接入 · AI 真的能帮你少踩坑",
        [
            "7 Markets × 3 Roles × 25 APIs = 40 Checks",
            "SG / MY / TH / VN / ID / PH / TW 一套覆盖",
            "Cursor · Claude Code · Codex · Gemini CLI 全平台",
        ],
        "要接一个支持七个市场的物流 API……你真的知道自己会踩几种坑吗……这一次……"
        "AI 真的能帮你……从选型、到代码、到 Webhook、到面单、到排障……一步不差。",
    ),
    (
        "Scene 2 · PAIN POINTS",
        "三类典型踩坑 · 文档第一页不会告诉你的事",
        [
            "签名：{app_id}_{timestamp}_{random}_{body} · 顺序错一个 = 10002",
            "Webhook：OrderCreate / WeightFee 走 1·5·15·60·180 分钟重试；Tracking 不重试",
            "面单 AWB：V1 按 tracking_no · V2 按 batch_no · 选错即 500 条限制",
        ],
        "同一个下单接口……七个东南亚市场字段规则完全不一样……HMAC-SHA256 签名四个字段拼错一次……"
        "就是 10002 签名失败……还要区分 secret_key 和 user_secret……Webhook 回调更刁钻……"
        "OrderCreate 和 WeightFee 失败后会按 1、5、15、60、180 分钟重试……但 Tracking 不会重试……"
        "面单到底该调 V1 按运单号拉……还是 V2 按批次号聚合拉……这些问题……没有一个是 API 文档第一页告诉你的。",
    ),
    (
        "Scene 3 · CAPABILITIES",
        "AI Agent Skill · 侧边栏里的物流接入专家",
        [
            "核心 6 能力：产品选型 · 示例代码 · 业务知识 · 质量评估 · 排障决策树 · Onboarding",
            "新增 2 扩展：Webhook 回调验签 · 面单 AWB 批量拉取",
            "纯插件形态：Cursor / Claude Code / Codex / Gemini CLI · 代码不落地",
        ],
        "我们把它做成了一个 AI Agent Skill 插件……装进 Cursor、Claude Code、Codex、Gemini CLI……"
        "所有能力都在侧边栏里……产品选型、示例代码、业务知识、质量评估、排障决策树、Onboarding 上线……"
        "再加上 Webhook 回调验签和面单 AWB 拉取……全链路都有专家级回答……而且所有代码只是参考范本……不会直接写进你的项目。",
    ),
    (
        "Scene 4 · QUALITY GATE",
        "代码质量 · 5 维度量化评分（B 级 78 分）",
        [
            "D1 签名与认证 25% · D2 订单流程 25% · D3 Webhook 处理 20%",
            "D4 错误处理 15% · D5 市场适配 15%",
            "Top Issues：硬编码 secret_key · 下单缺 try/catch · Webhook 未校验签名",
        ],
        "把一段真实的 PHP 接入代码丢进去……AI 扫出十五项具体问题……硬编码密钥、下单缺少异常处理、Webhook 回调没做签名校验……"
        "每一条都给出修复建议和对应 Playbook……五维度加权评分七十八分……刚好卡在 B 级建议修复……"
        "签名认证、订单流程、Webhook 处理、错误处理、市场适配……五个维度各自打分……你一眼就能看到……哪里该先补。",
    ),
    (
        "Scene 5 · CODE GEN + PRECHECK",
        "代码生成 · 参数预检 · 一套签名走到底",
        [
            "6 语言模板：Python · Java · Go · Node.js · PHP · cURL 全跑通",
            "Request + Webhook + AWB 共用同一套 HMAC-SHA256 签名逻辑",
            "JSON 预检：必填缺失 / 字段类型错 / 市场差异规则 一次性提示",
        ],
        "需要越南市场的示例代码吗……它给你完整可跑的模板……Python Java Go Node PHP cURL 六种语言全覆盖……"
        "关键是……Request 请求签名、Webhook 回调验签、AWB 面单签名……全部复用同一套 HMAC-SHA256 逻辑……"
        "再把真实 JSON 参数贴进来做预检……自动识别必填项缺失、字段格式错误、市场差异规则……一次性列清所有问题……不用上线再踩雷。",
    ),
    (
        "Scene 6 · ERROR ANALYSIS",
        "线上报错 → 根因 + 6 本 Playbook 一对一兜底",
        [
            "INVALID_PARAM / recipient.ward → 越南市场必填缺失 · 自动从地址簿补齐",
            "6 本 Playbook：签名 · 订单 · Webhook · 面单 · 运费 · 市场适配",
            "东南亚 7 市场高频场景全覆盖 · 不再在群里反复问",
        ],
        "线上报错 INVALID_PARAM 怎么办……AI 直接追到具体字段……给出根因分析……联动对应 Playbook 一对一兜底……"
        "签名、订单、Webhook、面单、运费、市场适配……六本专题手册覆盖东南亚七个市场的全部高频场景……再也不用在群里反复问同一个问题。",
    ),
    (
        "Scene 7 · CTA",
        "以前靠经验 · 现在靠评分 · Webhook + AWB 全链路覆盖",
        [
            "7 Markets · 40 Checks · 6 Playbooks · Webhook + AWB",
            "install spx-openapi-integration → Cursor / Claude Code / Codex / Gemini CLI",
            "告诉它：你是谁、去哪个市场 · 剩下的它替你走完",
        ],
        "以前对接 SPX 靠经验……现在可以靠评分……七个市场、四十项检查、六本 Playbook……"
        "Webhook 加 AWB 全链路覆盖……在 Cursor、Claude Code、Codex、Gemini CLI 里一键安装……"
        "你只告诉它你是谁、去哪个市场……剩下的……它替你走完。",
    ),
]


# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color: RGBColor) -> None:
    """给幻灯片填一个纯色底（托底，防止 PNG 边缘有透明时漏底）."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_centered_textbox(
    slide,
    text: str,
    *,
    top: Emu,
    height: Emu,
    font_size: int,
    font_color: RGBColor,
    bold: bool = False,
    letter_spacing: float = 0,
    prs: Presentation,
) -> None:
    """在幻灯片上添加一个横向居中、文本本身居中对齐的文本框."""
    width = prs.slide_width
    box = slide.shapes.add_textbox(Emu(0), top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    if letter_spacing:
        # python-pptx 对字间距支持有限，这里用制表符替代实现视觉上的间距
        pass


def add_caption_strip(
    slide,
    *,
    title: str,
    tag: str,
    prs: Presentation,
) -> None:
    """在底部 1/8 高度加一个半透明色带，左侧显示 tag、右侧显示标题."""
    strip_height = Inches(0.75)
    strip_top = prs.slide_height - strip_height
    box = slide.shapes.add_textbox(
        Emu(0), strip_top, prs.slide_width, strip_height,
    )
    # 给色带加一个半透明黑色填充
    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x07, 0x0A, 0x10)
    # line 去边
    line = box.line
    line.fill.background()

    tf = box.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    tag_run = p.add_run()
    tag_run.text = f"{tag}   "
    tag_run.font.size = Pt(14)
    tag_run.font.bold = True
    tag_run.font.color.rgb = ACCENT

    title_run = p.add_run()
    title_run.text = title
    title_run.font.size = Pt(20)
    title_run.font.bold = True
    title_run.font.color.rgb = TEXT


def add_notes(slide, notes_text: str) -> None:
    """把旁白写到备注区，方便演示时提词用."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    # 把 Unicode 三点省略号替换为普通省略号，避免某些 Office 下显示异常
    clean = notes_text.replace("……", "… ").replace("。", "。 ")
    run = p.add_run()
    run.text = clean
    run.font.size = Pt(14)


# ---------------------------------------------------------------------------
# 页面生成
# ---------------------------------------------------------------------------
def build_cover_slide(prs: Presentation) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG)

    add_centered_textbox(
        slide,
        "SPX OPEN API · AGENT SKILL",
        top=Inches(1.6),
        height=Inches(0.6),
        font_size=20,
        font_color=ACCENT,
        bold=True,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "物流 API 接入",
        top=Inches(2.2),
        height=Inches(1.4),
        font_size=72,
        font_color=TEXT,
        bold=True,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "这一次，AI 真的能帮你少踩坑",
        top=Inches(3.8),
        height=Inches(0.8),
        font_size=28,
        font_color=HIGHLIGHT,
        bold=True,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "7 Markets · 40 Checks · 6 Playbooks · Webhook + AWB",
        top=Inches(4.7),
        height=Inches(0.6),
        font_size=18,
        font_color=MUTED,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "Cursor   ·   Claude Code   ·   Codex   ·   Gemini CLI",
        top=Inches(5.4),
        height=Inches(0.6),
        font_size=16,
        font_color=GOLD,
        prs=prs,
    )

    add_notes(
        slide,
        "封面页。产品名：SPX OpenAPI 接入助手。一句话定位：让 AI 像资深物流接入顾问一样，"
        "帮你打分、查错、生成代码、定位根因。Webhook 与 AWB 全链路覆盖。"
        "以 AI Agent Skill 的形态装进 Cursor / Claude Code / Codex / Gemini CLI。",
    )


def build_scene_slide(
    prs: Presentation,
    *,
    image_path: Path,
    tag: str,
    title: str,
    key_points: List[str],
    narration: str,
) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG)

    # 图整页铺满
    slide.shapes.add_picture(
        str(image_path),
        Emu(0),
        Emu(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )

    # 底部标题色带
    add_caption_strip(slide, title=title, tag=tag, prs=prs)

    # 备注：要点 + 完整旁白
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = f"【{tag}】{title}"
    r0.font.bold = True
    r0.font.size = Pt(14)

    for point in key_points:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"· {point}"
        r.font.size = Pt(12)

    p_split = tf.add_paragraph()
    p_split.add_run().text = ""

    p_n = tf.add_paragraph()
    r_n = p_n.add_run()
    r_n.text = f"旁白：{narration}"
    r_n.font.size = Pt(12)


def build_closing_slide(prs: Presentation) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG)

    add_centered_textbox(
        slide,
        "以前靠经验 · 现在靠评分",
        top=Inches(1.4),
        height=Inches(1.2),
        font_size=56,
        font_color=TEXT,
        bold=True,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "AI 不是帮你写代码 · 是帮你少踩坑",
        top=Inches(2.7),
        height=Inches(0.8),
        font_size=28,
        font_color=HIGHLIGHT,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "install  spx-openapi-integration",
        top=Inches(3.8),
        height=Inches(0.8),
        font_size=28,
        font_color=GOLD,
        bold=True,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "Cursor   ·   Claude Code   ·   Codex   ·   Gemini CLI",
        top=Inches(4.6),
        height=Inches(0.6),
        font_size=18,
        font_color=ACCENT,
        prs=prs,
    )
    add_centered_textbox(
        slide,
        "#SPX · #物流API · #AgentSkill · #Webhook · #AWB · #东南亚物流",
        top=Inches(5.5),
        height=Inches(0.6),
        font_size=14,
        font_color=MUTED,
        prs=prs,
    )

    add_notes(
        slide,
        "结尾页。口号：以前靠经验，现在靠评分。CTA：install spx-openapi-integration。"
        "四大平台：Cursor / Claude Code / Codex / Gemini CLI。"
        "指标：7 Markets · 40 Checks · 6 Playbooks · Webhook + AWB。",
    )


# ---------------------------------------------------------------------------
# 入口
# ---------------------------------------------------------------------------
def main() -> None:
    if not SLIDES_DIR.exists():
        raise SystemExit(
            f"找不到静态帧目录：{SLIDES_DIR}\n"
            f"请先运行：npm run render:spxopenapi-landscape:slides",
        )

    # 16:9 标准尺寸（13.333 × 7.5 英寸，对应 1920×1080 的 PPT 画布）
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover_slide(prs)

    for idx, (tag, title, key_points, narration) in enumerate(SCENES, start=1):
        img_path = SLIDES_DIR / f"slide-{idx}.png"
        if not img_path.exists():
            raise SystemExit(
                f"找不到第 {idx} 页静态帧：{img_path}\n"
                f"请先运行：npm run render:spxopenapi-landscape:slides",
            )
        build_scene_slide(
            prs,
            image_path=img_path,
            tag=tag,
            title=title,
            key_points=key_points,
            narration=narration,
        )

    build_closing_slide(prs)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    print(f"[ok] 生成 PPT：{OUTPUT_PPTX}")
    print(f"[ok] 共 {len(prs.slides)} 页（1 封面 + {len(SCENES)} 场景 + 1 结尾）")


if __name__ == "__main__":
    main()
